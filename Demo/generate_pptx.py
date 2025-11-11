#!/usr/bin/env python3
"""
Generate TribalLink Presentation PPTX
Creates a professional 15-slide PowerPoint presentation for the college demo.
Run: python generate_pptx.py
Output: TribalLink_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
GREEN = RGBColor(76, 175, 80)      # #4CAF50
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(51, 51, 51)
ORANGE = RGBColor(255, 107, 53)    # #FF6B35

def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = GREEN
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = WHITE
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_list, bg_color=WHITE):
    """Add a content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = GREEN
    title_shape.line.color.rgb = GREEN
    
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.vertical_anchor = 1
    title_frame.margin_left = Inches(0.3)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)

# SLIDE 1: Title Slide
add_title_slide(prs, "TribalLink", "Smart Digital Network for Tribal Empowerment")

# SLIDE 2: The Problem
add_content_slide(prs, "🚨 The Challenge We're Solving", [
    "❌ Tribal artisans lack digital market access",
    "❌ Farmers have limited knowledge of modern techniques",
    "❌ Digital divide in rural communities",
    "❌ Limited knowledge sharing platforms",
    "❌ 104 million tribal people in India • Only 15% have internet access"
])

# SLIDE 3: Our Solution
add_content_slide(prs, "💡 Introducing TribalLink", [
    "✅ Digital marketplace for tribal products",
    "✅ AI-powered farming assistant (AgriHelp Bot)",
    "✅ Community connection platform",
    "✅ Mobile-responsive design",
    "✅ Easy to use (no technical skills needed)",
    "Tagline: Connecting Communities. Creating Opportunities. Changing Lives."
])

# SLIDE 4: Key Features
add_content_slide(prs, "✨ Platform Features", [
    "🏠 HOME PAGE - Welcome hub with project overview",
    "🛍️ MARKETPLACE - Browse & purchase tribal products (5+ categories)",
    "🤖 AGRIHELP BOT - AI chatbot for farming assistance",
    "📱 RESPONSIVE DESIGN - Works on smartphones, tablets, desktops",
    "Supports: Crafts, Jewelry, Pottery, Agriculture, Furniture"
])

# SLIDE 5: Technical Architecture
add_content_slide(prs, "🔧 How It Works - The Technology", [
    "FRONTEND: HTML5, CSS3, JavaScript (Responsive, Interactive)",
    "BACKEND: Python Flask, RESTful APIs, JSON Data Format",
    "FEATURES: Error handling, CORS, XSS prevention, Real-time chat",
    "DEPLOYMENT: Local Network (192.168.31.94:5000), Cloud-ready",
    "SCALABLE: Designed for millions of users"
])

# SLIDE 6: Live Demo
add_content_slide(prs, "🚀 Live Platform Demo", [
    "STEP 1: HOME PAGE - Welcome message, features, statistics",
    "STEP 2: MARKETPLACE - Browse products, apply filters (Crafts, Agriculture)",
    "STEP 3: CHATBOT - Ask 'How to grow rice?' → Get farming tips",
    "STEP 4: RESPONSIVENESS - Show mobile view (resize browser)",
    "Interactive, real-time responses powered by backend APIs"
])

# SLIDE 7: Team Members
add_content_slide(prs, "👥 Meet Team TechTribe", [
    "👨‍💼 HIMANSHU CHOUDHARY - Project Lead & Backend Developer",
    "👨‍🎨 SHEKHAR KUMAR - UI/UX Designer",
    "👨‍💻 AYUSH MANDAL - Documentation & Quality Assurance",
    "👨‍💼 ANKIT KUMAR GUPTA - Content & Feature Development",
    "Institution: Pemiya Rishikesh Institute of Technology • Competition: IDEA TRIBE 2025"
])

# SLIDE 8: Impact & Scalability
add_content_slide(prs, "📊 Impact & Growth Potential", [
    "CURRENT: 100+ artisans • 5+ categories • 24/7 AI support",
    "PHASE 2 (Q1 2026): 1,000+ artisans across 2-3 states",
    "PHASE 3 (Q2 2026): Mobile app, multi-language support",
    "PHASE 4 (Q3 2026): National expansion, advanced AI",
    "PHASE 5 (Q4 2026+): International reach • ₹10 crore annual volume"
])

# SLIDE 9: Social & Economic Impact
add_content_slide(prs, "🌱 Creating Real Change", [
    "💰 ECONOMIC: ₹5,000-10,000 monthly income increase per artisan",
    "👨‍👩‍👧‍👦 SOCIAL: Strengthen community, reduce brain drain to cities",
    "🌿 ENVIRONMENTAL: Sustainable farming, eco-friendly handicrafts",
    "✅ UN SDGs: No Poverty, Zero Hunger, Quality Education, Decent Work",
    "Empowering 10,000+ artisans and reaching 1M+ farmers"
])

# SLIDE 10: Competitive Advantages
add_content_slide(prs, "🎯 Why TribalLink Stands Out", [
    "1. COMMUNITY-CENTRIC: Built specifically for tribal communities",
    "2. INTEGRATED SOLUTION: Marketplace + Guidance + Community",
    "3. LOW TECH BARRIER: Works on any phone, poor internet OK",
    "4. SUSTAINABLE MODEL: Commission-based revenue (2-3%)",
    "5. GOVERNMENT READY: Aligns with rural development initiatives"
])

# SLIDE 11: Implementation Roadmap
add_content_slide(prs, "📅 Our Roadmap", [
    "✅ PHASE 1 (NOW): MVP complete, 100+ artisans, local launch",
    "📌 PHASE 2 (Q1 2026): Payment gateway, 1,000+ artisans",
    "📌 PHASE 3 (Q2 2026): Mobile app, multilingual support",
    "📌 PHASE 4 (Q3 2026): National expansion, advanced features",
    "📌 PHASE 5 (Q4 2026+): International expansion, investment round"
])

# SLIDE 12: Funding & Resources
add_content_slide(prs, "💼 Investment & Resources Needed", [
    "PHASE 1 (MVP): ₹5-10 lakhs ✓ COMPLETED",
    "PHASE 2-3 Budget: ₹1.25-2.5 crores",
    "USE OF FUNDS: Tech (40%), Team (30%), Marketing (20%), Operations (10%)",
    "REVENUE: Commission on sales (2-3%), Premium features, Training",
    "INVESTORS: Impact funds, Government grants, Corporate CSR"
])

# SLIDE 13: Challenges & Solutions
add_content_slide(prs, "⚠️ Challenges & Our Solutions", [
    "DIGITAL LITERACY → Simple interface + free training",
    "POOR INTERNET → Offline features + low-bandwidth design",
    "LIMITED BANKING → Multiple payment options (UPI, COD, bank)",
    "ADOPTION RESISTANCE → Community ambassadors + partnerships",
    "SUSTAINABILITY → Proven commission model + diversified revenue"
])

# SLIDE 14: Calls to Action
add_content_slide(prs, "🚀 How You Can Help", [
    "💰 FUNDING PARTNERS - Invest in sustainable impact",
    "🤝 TECHNOLOGY PARTNERS - Cloud infrastructure, payment gateway",
    "🌍 COMMUNITY PARTNERS - NGOs, tribal organizations, government",
    "👥 TALENT - Developers, designers, content creators (we're hiring!)",
    "📢 INDIVIDUALS - Share, use platform, refer artisans, give feedback"
])

# SLIDE 15: Closing
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = GREEN

# Closing title
closing_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
closing_frame = closing_box.text_frame
closing_frame.text = "🌍 Together, We Can Create Change"
closing_frame.paragraphs[0].font.size = Pt(44)
closing_frame.paragraphs[0].font.bold = True
closing_frame.paragraphs[0].font.color.rgb = WHITE
closing_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Vision
vision_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
vision_frame = vision_box.text_frame
vision_frame.text = "TribalLink: A digital ecosystem where tribal communities thrive.\nPreserving culture. Creating opportunity. Changing lives."
vision_frame.paragraphs[0].font.size = Pt(22)
vision_frame.paragraphs[0].font.color.rgb = WHITE
vision_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Save presentation
output_path = "TribalLink_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation created successfully: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
print("🎨 Color scheme: Green (#4CAF50) + Professional styling")
print("📝 Ready for college demo presentation!")
